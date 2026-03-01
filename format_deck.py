#!/usr/bin/env python3
"""
Phase 2: Comprehensive formatting and content polish for COBE Marketing Playbook.
Standardizes all fonts, colors, table styling, and adds professional polish.
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree

SRC = "COBE Marketing Playbook Updated (2).pptx"
DST = "COBE Marketing Playbook Updated (2).pptx"  # overwrite

# ─── Brand constants ──────────────────────────────────────────────────────────

BOISE_BLUE    = RGBColor(0x00, 0x33, 0xA0)
BOISE_ORANGE  = RGBColor(0xD6, 0x43, 0x09)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY     = RGBColor(0x33, 0x33, 0x33)
MED_GRAY      = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
BORDER_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)

# Funnel tag colors (brand-aligned)
TAG_AWARENESS    = BOISE_BLUE
TAG_CONSIDERATION = BOISE_ORANGE
TAG_CONVERSION   = RGBColor(0x1B, 0x7A, 0x2F)  # Professional green

# Activation box backgrounds
HANDS_OFF_BG = RGBColor(0xE8, 0xED, 0xF7)  # Light blue tint
HANDS_ON_BG  = RGBColor(0xFD, 0xF0, 0xE6)  # Light orange tint

# Deck font
FONT = 'Arial'

# Standard positions (EMU)
LEFT_MARGIN = Emu(457200)    # 0.5"
RIGHT_EDGE  = Emu(8686800)   # 9.5"
CONTENT_W   = Emu(8229600)   # 9.0"
FOOTER_Y    = Emu(5000000)   # near bottom
FOOTER_H    = Emu(143500)

SLIDE_W = 9144000
SLIDE_H = 5143500

# ─── Helpers ──────────────────────────────────────────────────────────────────

def is_new_shape(shape):
    """Return True if shape was added by the previous script (not from Google Slides)."""
    name = shape.name
    return not name.startswith('Google Shape')

def standardize_run_font(run, override_name=None, override_size=None,
                         override_bold=None, override_color=None):
    """Set font properties on a run. Only sets non-None values."""
    if override_name is not None:
        run.font.name = override_name
    elif run.font.name and run.font.name != FONT:
        run.font.name = FONT
    if override_size is not None:
        run.font.size = override_size
    if override_bold is not None:
        run.font.bold = override_bold
    if override_color is not None:
        run.font.color.rgb = override_color

def fix_all_fonts_to_arial(shape):
    """Fix all runs in a shape to use Arial."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name != FONT:
                run.font.name = FONT

def fix_color_999_to_666(shape):
    """Fix #999999 colors to #666666."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                if run.font.color and run.font.color.rgb == RGBColor(0x99, 0x99, 0x99):
                    run.font.color.rgb = MED_GRAY
            except:
                pass

def remove_new_shapes(slide):
    """Remove all shapes that were added by the previous script."""
    to_remove = []
    for shape in slide.shapes:
        if is_new_shape(shape):
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    return len(to_remove)

def add_textbox(slide, left, top, width, height, text,
                font_size=10, bold=False, color=BLACK,
                alignment=PP_ALIGN.LEFT, font_name=FONT,
                word_wrap=True, v_anchor=None):
    """Add a text box with single-line text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if v_anchor:
        tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_rich_textbox(slide, left, top, width, height, lines):
    """Add textbox with multiple paragraphs. Each line is a dict:
    {text, size, bold, color, alignment, space_after, space_before}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, cfg in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = cfg.get('alignment', PP_ALIGN.LEFT)
        p.space_after = Pt(cfg.get('space_after', 2))
        p.space_before = Pt(cfg.get('space_before', 0))
        run = p.add_run()
        run.text = cfg['text']
        run.font.size = Pt(cfg.get('size', 10))
        run.font.bold = cfg.get('bold', False)
        run.font.color.rgb = cfg.get('color', DARK_GRAY)
        run.font.name = FONT
    return txBox

def add_funnel_tag(slide, left, top, text, bg_color):
    """Add a professional funnel indicator tag."""
    # Calculate width based on text length
    char_width = Emu(60000)
    width = max(Emu(600000), len(text) * char_width + Emu(180000))
    height = Emu(228600)  # ~0.25"

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.15

    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT
    return shape, width

def add_activation_section(slide, left, top, width,
                           ho_items, hn_items, channel_name=""):
    """Add professional activation section matching deck visual language."""
    # Section header
    hdr = add_textbox(slide, left, top, width, Emu(200000),
                      "HOW TO ACTIVATE", font_size=10, bold=True,
                      color=BOISE_ORANGE)

    # Calculate box dimensions
    box_w = width
    box_h_ho = Emu(700000)
    box_h_hn = Emu(700000)
    gap = Emu(70000)

    # ── Hands-Off Box ──
    ho_y = top + Emu(220000)
    ho_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, ho_y, box_w, box_h_ho
    )
    ho_box.fill.solid()
    ho_box.fill.fore_color.rgb = HANDS_OFF_BG
    ho_box.line.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
    ho_box.line.width = Pt(0.75)
    ho_box.adjustments[0] = 0.08

    tf = ho_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(5)
    tf.margin_bottom = Pt(3)

    # Header
    p = tf.paragraphs[0]
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = "HANDS-OFF (Concierge)"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = BOISE_BLUE
    run.font.name = FONT

    # Email instruction
    p2 = tf.add_paragraph()
    p2.space_after = Pt(1)
    run2 = p2.add_run()
    run2.text = "Email COBE Marketing at cobe-info@boisestate.edu with:"
    run2.font.size = Pt(8)
    run2.font.color.rgb = DARK_GRAY
    run2.font.name = FONT

    # Items
    for item in ho_items:
        pi = tf.add_paragraph()
        pi.space_after = Pt(0)
        ri = pi.add_run()
        ri.text = item
        ri.font.size = Pt(7.5)
        ri.font.color.rgb = DARK_GRAY
        ri.font.name = FONT

    # ── Hands-On Box ──
    hn_y = ho_y + box_h_ho + gap
    hn_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, hn_y, box_w, box_h_hn
    )
    hn_box.fill.solid()
    hn_box.fill.fore_color.rgb = HANDS_ON_BG
    hn_box.line.color.rgb = RGBColor(0xED, 0xC9, 0xAF)
    hn_box.line.width = Pt(0.75)
    hn_box.adjustments[0] = 0.08

    tf2 = hn_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Pt(8)
    tf2.margin_right = Pt(8)
    tf2.margin_top = Pt(5)
    tf2.margin_bottom = Pt(3)

    p = tf2.paragraphs[0]
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = "HANDS-ON (DIY w/ Support)"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = BOISE_ORANGE
    run.font.name = FONT

    p2 = tf2.add_paragraph()
    p2.space_after = Pt(1)
    run2 = p2.add_run()
    run2.text = "Draft using COBE templates \u2192 email cobe-info@boisestate.edu for review/approval:"
    run2.font.size = Pt(8)
    run2.font.color.rgb = DARK_GRAY
    run2.font.name = FONT

    for item in hn_items:
        pi = tf2.add_paragraph()
        pi.space_after = Pt(0)
        ri = pi.add_run()
        ri.text = item
        ri.font.size = Pt(7.5)
        ri.font.color.rgb = DARK_GRAY
        ri.font.name = FONT

def add_screenshot_frame(slide, left, top, width, height, label, sublabel=""):
    """Add a professional screenshot placeholder frame."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    shape.line.color.rgb = BOISE_BLUE
    shape.line.width = Pt(1.5)
    # Dashed border
    from pptx.oxml.ns import qn
    ln = shape._element.find(qn('a:ln'), shape._element.nsmap) or shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    if ln is not None:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(4)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = BOISE_BLUE
    run.font.name = FONT

    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sublabel
        r2.font.size = Pt(7)
        r2.font.color.rgb = MED_GRAY
        r2.font.name = FONT
        r2.font.italic = True

    return shape

def add_slide_footer(slide, slide_num, total_slides):
    """Add a professional footer with slide number."""
    # Right-aligned page number
    num_box = add_textbox(
        slide, Emu(7800000), FOOTER_Y, Emu(886800), FOOTER_H,
        f"{slide_num}", font_size=7, color=MED_GRAY,
        alignment=PP_ALIGN.RIGHT
    )
    # Left-aligned branding
    brand_box = add_textbox(
        slide, LEFT_MARGIN, FOOTER_Y, Emu(4000000), FOOTER_H,
        "COBE Marketing Playbook  |  cobe-info@boisestate.edu",
        font_size=7, color=MED_GRAY, alignment=PP_ALIGN.LEFT
    )
    return num_box, brand_box

def format_table_cells(table, header_bg=BOISE_BLUE, header_fg=WHITE,
                       alt_row_bg=LIGHT_GRAY, body_fg=DARK_GRAY):
    """Apply consistent formatting to a table."""
    for ri, row in enumerate(table.rows):
        for ci, cell in enumerate(row.cells):
            # Background
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_row_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

            # Cell margins
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)

            # Text formatting
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT
                    if ri == 0:
                        run.font.size = Pt(9)
                        run.font.bold = True
                        run.font.color.rgb = header_fg
                    else:
                        run.font.size = Pt(8)
                        run.font.bold = False
                        run.font.color.rgb = body_fg

# ─── Slide-specific formatting functions ──────────────────────────────────────

def format_channel_slide(slide, slide_idx, channel_name, funnel_tags,
                         ho_items, hn_items, screenshot_label, screenshot_sub):
    """Full format pass for a channel slide: remove old new-shapes, add fresh ones."""
    # Step 1: Remove previously-added shapes
    removed = remove_new_shapes(slide)
    print(f"    Removed {removed} old new-shapes")

    # Step 2: Add funnel tags (top-right, stacked horizontally)
    tag_y = Emu(274320)
    tag_x = Emu(5943600)

    tag_colors = {
        "Awareness": TAG_AWARENESS,
        "Consideration": TAG_CONSIDERATION,
        "Conversion": TAG_CONVERSION,
    }

    x_cursor = tag_x
    for tag_name in funnel_tags:
        color = tag_colors.get(tag_name, BOISE_BLUE)
        _, w = add_funnel_tag(slide, x_cursor, tag_y, tag_name, color)
        x_cursor += w + Emu(70000)

    # Step 3: Add activation section (right side, below existing content)
    add_activation_section(
        slide,
        left=Emu(5029200),
        top=Emu(3150000),
        width=Emu(3657600),
        ho_items=ho_items,
        hn_items=hn_items,
        channel_name=channel_name,
    )

    # Step 4: Add screenshot placeholder
    add_screenshot_frame(
        slide,
        Emu(457200), Emu(3850000), Emu(4400000), Emu(800000),
        screenshot_label, screenshot_sub
    )

    print(f"    Added funnel tags: {funnel_tags}")
    print(f"    Added activation boxes + screenshot placeholder")


def rebuild_templates_hub(slide):
    """Completely rebuild Templates & Links Hub with deck-consistent formatting."""
    # Remove ALL shapes
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # ── Title ──
    add_textbox(slide, LEFT_MARGIN, Emu(274320), CONTENT_W, Emu(548640),
                "Templates & Links Hub", font_size=32, bold=True,
                color=BOISE_BLUE)

    # ── Accent bar ──
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, LEFT_MARGIN, Emu(777240), Emu(1371600), Emu(54864)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BOISE_ORANGE
    bar.line.fill.background()

    # ── Subtitle ──
    add_textbox(slide, LEFT_MARGIN, Emu(914400), CONTENT_W, Emu(274320),
                "Everything you need to get started \u2014 templates, tools, and reference links",
                font_size=12, color=MED_GRAY)

    # ── Template cards ──
    templates = [
        ("Creative Brief Template",
         "Submit campaign details to COBE Marketing",
         "templates/creative-brief-template.md",
         "Use for every new campaign or event request"),
        ("Copy Prompts Template",
         "Pre-written copy frameworks for each channel",
         "templates/copy-prompts-template.md",
         "Instagram, LinkedIn, email, screens, classroom"),
        ("Campaign Calendar Template",
         "Plan multi-channel event promotion timelines",
         "templates/campaign-calendar-template.md",
         "Work-back schedule: 6 weeks to event day"),
        ("Budget Request Template",
         "Estimate costs and request paid advertising",
         "templates/budget-request-template.md",
         "CPM benchmarks + free vs paid channel mix"),
        ("Social Media Calendar",
         "Weekly/monthly posting schedule",
         "Social Media Calendar Template.xlsx",
         "Excel template for content planning"),
        ("Social Media Calendar Example",
         "Sample calendar showing structure and pacing",
         "Social Media Calendar Example.xlsx",
         "Model for cadence and content types"),
    ]

    col_w = Emu(3931920)
    row_h = Emu(550000)
    start_y = Emu(1280000)
    start_x = LEFT_MARGIN
    gap_x = Emu(365760)
    gap_y = Emu(100000)

    for idx, (name, desc, path, detail) in enumerate(templates):
        col = idx % 2
        row = idx // 2
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, row_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = BORDER_GRAY
        card.line.width = Pt(0.75)
        card.adjustments[0] = 0.06

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(8)
        tf.margin_bottom = Pt(4)

        # Template name
        p1 = tf.paragraphs[0]
        p1.space_after = Pt(2)
        r1 = p1.add_run()
        r1.text = name
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = BOISE_BLUE
        r1.font.name = FONT

        # Description
        p2 = tf.add_paragraph()
        p2.space_after = Pt(1)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(9)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = FONT

        # Detail
        p3 = tf.add_paragraph()
        p3.space_after = Pt(1)
        r3 = p3.add_run()
        r3.text = detail
        r3.font.size = Pt(8)
        r3.font.color.rgb = MED_GRAY
        r3.font.name = FONT
        r3.font.italic = True

        # File path
        p4 = tf.add_paragraph()
        r4 = p4.add_run()
        r4.text = f"\u2192 {path}"
        r4.font.size = Pt(7)
        r4.font.color.rgb = BOISE_BLUE
        r4.font.name = FONT

    # ── Key Links box at bottom ──
    links_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Emu(4200000), CONTENT_W, Emu(700000)
    )
    links_bg.fill.solid()
    links_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xED, 0xF7)
    links_bg.line.fill.background()
    links_bg.adjustments[0] = 0.04

    add_rich_textbox(slide, Emu(594360), Emu(4230000), Emu(7955280), Emu(650000), [
        {"text": "KEY LINKS & CONTACTS", "size": 10, "bold": True,
         "color": BOISE_BLUE, "space_after": 3},
        {"text": "All requests: cobe-info@boisestate.edu   |   Subject line: PROGRAM REQUEST - Channel - Date Range",
         "size": 8, "color": DARK_GRAY, "space_after": 2},
        {"text": "University events calendar: boisestate.enterprise.localist.com   |   Brand guide: brandguide.boisestate.edu",
         "size": 8, "color": DARK_GRAY, "space_after": 2},
        {"text": "Career Services: cobecareers@boisestate.edu   |   QR code generator: flowcode.com",
         "size": 8, "color": DARK_GRAY, "space_after": 0},
    ])


def rebuild_faq_slide(slide, templates_hub_num):
    """Polish FAQ slide with correct cross-references and formatting."""
    # Remove only the new shape we added (FAQ template question)
    remove_new_shapes(slide)

    # Add new FAQ entries at the bottom
    add_rich_textbox(slide, LEFT_MARGIN, Emu(4750000), CONTENT_W, Emu(350000), [
        {"text": f"Q: Where do I find the templates?", "size": 11, "bold": True,
         "color": DARK_GRAY, "space_after": 1},
        {"text": f"A: See the Templates & Links Hub (Slide {templates_hub_num}) or the /templates folder in the playbook repository.",
         "size": 11, "bold": False, "color": MED_GRAY, "space_after": 0},
    ])


def format_slide10_funnel(slide):
    """Polish the Marketing Funnel slide's 'How to use this' callout."""
    remove_new_shapes(slide)

    # Add polished callout box
    callout_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Emu(4650000), CONTENT_W, Emu(420000)
    )
    callout_bg.fill.solid()
    callout_bg.fill.fore_color.rgb = HANDS_OFF_BG
    callout_bg.line.fill.background()
    callout_bg.adjustments[0] = 0.06

    add_rich_textbox(slide, Emu(594360), Emu(4680000), Emu(7955280), Emu(360000), [
        {"text": "HOW TO USE THIS FUNNEL:", "size": 9, "bold": True,
         "color": BOISE_BLUE, "space_after": 2},
        {"text": "1. Define your goal (awareness / signups / attendance)     "
                 "2. Choose the funnel stage     "
                 "3. Go to the recommended channel page for activation details, specs, and lead times",
         "size": 8, "color": DARK_GRAY, "space_after": 0},
    ])


def format_internal_channels(slide):
    """Polish Internal Channels slide."""
    remove_new_shapes(slide)

    # Funnel tags
    tag_y = Emu(274320)
    x = Emu(6300000)
    _, w1 = add_funnel_tag(slide, x, tag_y, "Awareness", TAG_AWARENESS)
    add_funnel_tag(slide, x + w1 + Emu(70000), tag_y, "Conversion", TAG_CONVERSION)

    # Activation boxes
    add_activation_section(
        slide,
        left=Emu(4754880), top=Emu(3200000), width=Emu(3931920),
        ho_items=[
            "\u2022 Event details, dates, QR code destination URL",
            "\u2022 Photos or graphics (if available)",
            "\u2022 Target audience and placement preference",
        ],
        hn_items=[
            "\u2022 Finished slide (1920x1080, dark background, no logo)",
            "\u2022 QR code linking to registration/RSVP page",
        ],
    )

    # Lead time + specs footer
    specs_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Emu(4600000), CONTENT_W, Emu(350000)
    )
    specs_bg.fill.solid()
    specs_bg.fill.fore_color.rgb = LIGHT_GRAY
    specs_bg.line.fill.background()
    specs_bg.adjustments[0] = 0.06

    add_rich_textbox(slide, Emu(594360), Emu(4620000), Emu(7955280), Emu(300000), [
        {"text": "LEAD TIME & SPECS", "size": 8, "bold": True, "color": BOISE_ORANGE, "space_after": 2},
        {"text": "Classroom slides: 5 business days   |   MBEB screens: 5 business days   |   "
                 "1920x1080px, dark background, no logo overlays, large QR codes for scanning",
         "size": 8, "color": DARK_GRAY, "space_after": 0},
    ])


def format_email_slide(slide):
    """Polish Email Marketing channel slide."""
    remove_new_shapes(slide)

    # Funnel tag
    add_funnel_tag(slide, Emu(7400000), Emu(274320), "Conversion", TAG_CONVERSION)

    # Activation boxes (left-bottom area, below the rules)
    add_activation_section(
        slide,
        left=LEFT_MARGIN, top=Emu(3100000), width=Emu(3800000),
        ho_items=[
            "\u2022 Event/program details + target audience",
            "\u2022 Copy (\u2264100 words per item), CTA, final link",
            "\u2022 Images with alt text descriptions",
        ],
        hn_items=[
            "\u2022 Full copy draft + images",
            "\u2022 Target segment (students / alumni / employers)",
            "\u2022 Subject line (\u226455 characters)",
        ],
    )

    # Screenshot placeholder (right-bottom)
    add_screenshot_frame(
        slide,
        Emu(4754880), Emu(3100000), Emu(3931920), Emu(780000),
        "INSERT SCREENSHOT",
        "Replace with newsletter section example"
    )

    # Rename newsletter strategy box to include lead time
    # (already done in phase 1, just verify the content box exists)


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(SRC)
    slides = list(prs.slides)
    total = len(slides)

    print(f"Starting comprehensive formatting polish...")
    print(f"Total slides: {total}\n")

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 1: Global font and color standardization
    # ═══════════════════════════════════════════════════════════════════════════
    print("Phase 1: Global font & color standardization")
    font_fixes = 0
    color_fixes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and run.font.name != FONT:
                            run.font.name = FONT
                            font_fixes += 1
                        try:
                            if run.font.color and run.font.color.rgb == RGBColor(0x99, 0x99, 0x99):
                                run.font.color.rgb = MED_GRAY
                                color_fixes += 1
                        except:
                            pass
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name and run.font.name != FONT:
                                    run.font.name = FONT
                                    font_fixes += 1
    print(f"  Fixed {font_fixes} font names to Arial")
    print(f"  Fixed {color_fixes} #999 colors to #666")

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 2: Table formatting
    # ═══════════════════════════════════════════════════════════════════════════
    print("\nPhase 2: Table formatting")
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_table:
                print(f"  Formatting table on slide {i+1}")
                format_table_cells(shape.table)

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 3: Channel slide polish (remove old + rebuild)
    # ═══════════════════════════════════════════════════════════════════════════
    print("\nPhase 3: Channel slide polish")

    # Find channel slides by their label shapes
    channel_map = {}
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t == "CHANNEL":
                    # Determine which channel by looking at the title shape
                    for s2 in slide.shapes:
                        if s2.has_text_frame:
                            t2 = s2.text_frame.text.strip()
                            if t2 in ["Instagram", "LinkedIn", "Email Marketing"]:
                                channel_map[t2] = (i, slide)
                                break
                elif t == "INTERNAL":
                    channel_map["Internal Channels"] = (i, slide)

    # Instagram
    if "Instagram" in channel_map:
        idx, slide = channel_map["Instagram"]
        print(f"  Slide {idx+1}: Instagram")
        format_channel_slide(
            slide, idx, "Instagram",
            funnel_tags=["Awareness", "Consideration"],
            ho_items=[
                "\u2022 Event/program details (who, what, when, where)",
                "\u2022 Photos (minimum 1080px resolution)",
                "\u2022 Registration or RSVP link",
                "\u2022 Goal/success metric (e.g., 50 RSVPs)",
            ],
            hn_items=[
                "\u2022 Finished graphic (1080x1080 or 1080x1920)",
                "\u2022 Caption text with clear CTA",
                "\u2022 Relevant hashtags",
            ],
            screenshot_label="INSERT SCREENSHOT",
            screenshot_sub="Replace with Instagram feed or story example"
        )

    # LinkedIn
    if "LinkedIn" in channel_map:
        idx, slide = channel_map["LinkedIn"]
        print(f"  Slide {idx+1}: LinkedIn")
        format_channel_slide(
            slide, idx, "LinkedIn",
            funnel_tags=["Consideration", "Conversion"],
            ho_items=[
                "\u2022 Story/achievement details + outcome data",
                "\u2022 Partner names to tag",
                "\u2022 Photos, headshots, or PDF carousel",
                "\u2022 Registration or event link",
            ],
            hn_items=[
                "\u2022 Finished post copy + graphic (1920x1080)",
                "\u2022 Partner tags and relevant links",
            ],
            screenshot_label="INSERT SCREENSHOT",
            screenshot_sub="Replace with LinkedIn post or article example"
        )

    # Email Marketing
    if "Email Marketing" in channel_map:
        idx, slide = channel_map["Email Marketing"]
        print(f"  Slide {idx+1}: Email Marketing")
        format_email_slide(slide)

    # Internal Channels
    if "Internal Channels" in channel_map:
        idx, slide = channel_map["Internal Channels"]
        print(f"  Slide {idx+1}: Internal Channels")
        format_internal_channels(slide)

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 4: Marketing Funnel callout polish
    # ═══════════════════════════════════════════════════════════════════════════
    print("\nPhase 4: Marketing Funnel polish")
    # Find slide 10 (Marketing Funnel)
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "The Marketing Funnel":
                print(f"  Slide {i+1}: Marketing Funnel")
                format_slide10_funnel(slide)
                break

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 5: Templates & Links Hub rebuild
    # ═══════════════════════════════════════════════════════════════════════════
    print("\nPhase 5: Templates & Links Hub rebuild")
    templates_hub_num = None
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "Templates & Links Hub":
                templates_hub_num = i + 1
                print(f"  Slide {i+1}: Rebuilding Templates & Links Hub")
                rebuild_templates_hub(slide)
                break
        if templates_hub_num:
            break

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 6: FAQ polish
    # ═══════════════════════════════════════════════════════════════════════════
    print("\nPhase 6: FAQ polish")
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "Frequently Asked Questions":
                print(f"  Slide {i+1}: FAQ - updating with cross-references")
                rebuild_faq_slide(slide, templates_hub_num or "33")
                break

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 7: Slide footers (page numbers + branding)
    # ═══════════════════════════════════════════════════════════════════════════
    print("\nPhase 7: Adding slide footers")
    # Identify section divider slides (dark background) - skip footers on those
    section_dividers = set()
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t.startswith("SECTION ") or t == "PROOF OF CONCEPT":
                    section_dividers.add(i)
                    break
    # Also skip title slide (0) and closing slide (last)
    skip_slides = section_dividers | {0, total - 1}

    footer_count = 0
    for i, slide in enumerate(slides):
        if i not in skip_slides:
            add_slide_footer(slide, i + 1, total)
            footer_count += 1
    print(f"  Added footers to {footer_count} content slides")
    print(f"  Skipped {len(skip_slides)} slides (title, dividers, closing)")

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 8: Content refinements
    # ═══════════════════════════════════════════════════════════════════════════
    print("\nPhase 8: Content refinements")

    # Fix the FAQ "how fast" question to avoid "you" directed at COBE Marketing
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "How fast can you turn this around?" in shape.text_frame.text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "How fast can you turn this around?" in run.text:
                            run.text = run.text.replace(
                                "How fast can you turn this around?",
                                "How fast can COBE Marketing turn this around?")
                            print(f"  Fixed FAQ question wording on slide {i+1}")

    # Fix the Slide 18 title to have proper size (was smaller than other titles)
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "Sample Request Email" in shape.text_frame.text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size < Pt(32):
                            run.font.size = Pt(32)
                            print(f"  Fixed title font size on slide {i+1}")
                break

    # Add "Hands-Off" label to the sample email footer note
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if "Include: Event details, channels requested" in t:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Include:" in run.text:
                                run.text = "HANDS-OFF REQUEST TIP: Include event details, channels requested, goal/success metric, assets, and any links."
                                print(f"  Updated sample email footer on slide {i+1}")

    # Ensure Quick-Start slide tip about funnel references correct slide number
    funnel_slide_num = None
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "The Marketing Funnel" in shape.text_frame.text:
                funnel_slide_num = i + 1
                break
        if funnel_slide_num:
            break

    if funnel_slide_num:
        for i, slide in enumerate(slides):
            for shape in slide.shapes:
                if shape.has_text_frame and "see Marketing Funnel slide" in shape.text_frame.text:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "see Marketing Funnel slide" in run.text:
                                run.text = run.text.replace(
                                    "see Marketing Funnel slide",
                                    f"see Marketing Funnel, Slide {funnel_slide_num}")
                                print(f"  Updated funnel cross-reference on slide {i+1}")

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 9: Final global font pass (catch anything missed)
    # ═══════════════════════════════════════════════════════════════════════════
    print("\nPhase 9: Final font consistency pass")
    final_fixes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and run.font.name != FONT:
                            run.font.name = FONT
                            final_fixes += 1
    print(f"  Fixed {final_fixes} additional font inconsistencies")

    # ═══════════════════════════════════════════════════════════════════════════
    # Save
    # ═══════════════════════════════════════════════════════════════════════════
    print(f"\nSaving to {DST}...")
    prs.save(DST)
    print(f"Done! Final slide count: {len(prs.slides)}")

if __name__ == "__main__":
    main()
