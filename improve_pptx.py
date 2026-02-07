#!/usr/bin/env python3
"""
COBE Marketing Playbook — Improvement Script
Applies design, formatting, structural, and content improvements
to the existing PPTX while preserving all original content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ─── CONSTANTS ────────────────────────────────────────────────────────────────

INPUT_FILE = "COBE Marketing Playbook Updated After Comments.pptx"
OUTPUT_FILE = "COBE_Marketing_Playbook (3) — Improved.pptx"

# Boise State brand colours
BSU_BLUE   = RGBColor(0x00, 0x33, 0xA0)
BSU_ORANGE = RGBColor(0xD6, 0x43, 0x09)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Golden Rule - canonical format
GOLDEN_EMAIL = "cobe-info@boisestate.edu"
GOLDEN_SUBJECT = "PROGRAM REQUEST \u2013 [Channel] \u2013 [Date Range]"

# Font standards
FONT_FAMILY = "Arial"


# ─── HELPERS ──────────────────────────────────────────────────────────────────

def set_font(run, name=FONT_FAMILY, size=None, bold=None, italic=None,
             color=None, underline=None):
    """Apply font properties to a run."""
    run.font.name = name
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    if underline is not None:
        run.font.underline = underline


def add_textbox(slide, left, top, width, height, text="",
                font_size=12, bold=False, color=BLACK,
                alignment=PP_ALIGN.LEFT, font_name=FONT_FAMILY):
    """Add a simple text box and return the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, name=font_name, size=font_size, bold=bold, color=color)
    return txBox


def add_rounded_rect(slide, left, top, width, height,
                     fill_color=None, line_color=None):
    """Add a rounded rectangle shape and return it."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rectangle(slide, left, top, width, height,
                  fill_color=None, line_color=None):
    """Add a rectangle shape and return it."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_para(text_frame, text, font_size=12, bold=False, color=BLACK,
             alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(2),
             font_name=FONT_FAMILY, hyperlink=None):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    set_font(run, name=font_name, size=font_size, bold=bold, color=color)
    if hyperlink:
        run.hyperlink.address = hyperlink
    return p


def add_hyperlink_run(paragraph, text, url, font_size=12, bold=False,
                      color=BSU_BLUE, font_name=FONT_FAMILY, underline=True):
    """Add a hyperlinked run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    set_font(run, name=font_name, size=font_size, bold=bold, color=color,
             underline=underline)
    run.hyperlink.address = url
    return run


def add_footer(slide, slide_num):
    """Add consistent footer to a slide."""
    # Footer text
    add_textbox(slide, Inches(0.5), Inches(5.25), Inches(6), Inches(0.3),
                text=f"COBE Marketing Playbook  |  {GOLDEN_EMAIL}",
                font_size=8, color=MED_GRAY, alignment=PP_ALIGN.LEFT)
    # Slide number
    add_textbox(slide, Inches(8.5), Inches(5.25), Inches(1), Inches(0.3),
                text=str(slide_num),
                font_size=8, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


def add_path_tag(slide, path_num, label):
    """Add a PATH TAG at top-right of a slide."""
    tag_w = Inches(2.2)
    tag_h = Inches(0.3)
    tag_left = SLIDE_W - tag_w - Inches(0.15)
    tag_top = Inches(0.1)

    if path_num == 1:
        fill = BSU_BLUE
    else:
        fill = BSU_ORANGE

    tag = add_rounded_rect(slide, tag_left, tag_top, tag_w, tag_h,
                           fill_color=fill)
    tf = tag.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_font(run, size=9, bold=True, color=WHITE)
    # Vertically center
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    tag.text_frame.margin_top = Pt(2)
    tag.text_frame.margin_bottom = Pt(2)


def make_button_shape(slide, left, top, width, height, label, url,
                      fill_color=BSU_BLUE, font_size=10):
    """Create a clickable button shape with hyperlink."""
    btn = add_rounded_rect(slide, left, top, width, height,
                           fill_color=fill_color)
    tf = btn.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_font(run, size=font_size, bold=True, color=WHITE)
    if url:
        run.hyperlink.address = url
    return btn


def insert_slide_at(prs, slide, position):
    """Move a slide to a specific position (0-indexed) in the presentation."""
    from pptx.oxml.ns import qn as _qn
    # Get the slide ID list from the presentation XML element
    sldIdLst = prs.element.find(_qn('p:sldIdLst'))
    sldId_elements = list(sldIdLst)

    # The new slide is at the end
    new_sldId = sldId_elements[-1]

    # Remove it from current position
    sldIdLst.remove(new_sldId)

    # Re-read elements after removal
    remaining = list(sldIdLst)

    # Insert at the desired position
    if position == 0 or len(remaining) == 0:
        sldIdLst.insert(0, new_sldId)
    elif position >= len(remaining):
        sldIdLst.append(new_sldId)
    else:
        remaining[position].addprevious(new_sldId)


def add_blank_slide(prs):
    """Add a slide using the Blank layout."""
    blank_layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(blank_layout)


# ─── MAIN IMPROVEMENT PIPELINE ───────────────────────────────────────────────

def main():
    prs = Presentation(INPUT_FILE)
    slides = list(prs.slides)
    total_original = len(slides)
    print(f"Loaded {total_original} slides from {INPUT_FILE}")

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 1: PATH TAGS on existing slides
    # ═══════════════════════════════════════════════════════════════════════
    print("Phase 1: Adding PATH TAGS...")

    # PATH 1 slides (0-indexed): 11-22 (slides 12-23 in deck numbering)
    # Slide 11 (idx 10) is divider — skip tag on dividers
    path1_content_indices = list(range(11, 23))  # slides 12-23 (0-indexed 11-22)
    for idx in path1_content_indices:
        if idx < len(slides):
            add_path_tag(slides[idx], 1, "PATH 1: HANDS-OFF")

    # PATH 2 slides (0-indexed): 24-34 (slides 25-35 in deck numbering)
    # Slide 24 (idx 23) is divider — skip tag on dividers
    path2_content_indices = list(range(24, 35))  # slides 25-35 (0-indexed 24-34)
    for idx in path2_content_indices:
        if idx < len(slides):
            add_path_tag(slides[idx], 2, "PATH 2: HANDS-ON")

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 2: Golden Rule consistency — fix subject format dashes
    # ═══════════════════════════════════════════════════════════════════════
    print("Phase 2: Golden Rule consistency pass...")

    # We use en dash \u2013 consistently
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Fix subject line format: replace hyphen-based patterns
                        # with en-dash in "PROGRAM REQUEST" contexts
                        if "PROGRAM REQUEST" in run.text or "REQUEST" in run.text:
                            if " - " in run.text and "\u2013" not in run.text:
                                run.text = run.text.replace(" - ", " \u2013 ")

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 3: Add clickable hyperlinks to existing referenced URLs
    # ═══════════════════════════════════════════════════════════════════════
    print("Phase 3: Adding clickable hyperlinks to existing content...")

    # Slide 30 (idx 29) — WordPress accessibility training URL
    slide_30 = slides[29]
    for shape in slide_30.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "boisestate.edu/webguide/wordpress-accessibility" in run.text:
                        run.hyperlink.address = "https://www.boisestate.edu/webguide/wordpress-accessibility-compliance-training/"
                    if "helpdesk@boisestate.edu" in run.text:
                        run.hyperlink.address = "mailto:helpdesk@boisestate.edu"

    # Slide 32 (idx 31) — Flowcode
    slide_32 = slides[31]
    for shape in slide_32.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Flowcode" in run.text and "free" in run.text.lower():
                        run.hyperlink.address = "https://www.flowcode.com"

    # Slide 40 (idx 39) — Key links
    slide_40 = slides[39]
    for shape in slide_40.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "flowcode.com" in run.text.lower():
                        run.hyperlink.address = "https://www.flowcode.com"
                    if "brandguide.boisestate.edu" in run.text.lower():
                        run.hyperlink.address = "https://brandguide.boisestate.edu"
                    if "cobecareers@boisestate.edu" in run.text:
                        run.hyperlink.address = "mailto:cobecareers@boisestate.edu"
                    if "cobe-info@boisestate.edu" in run.text and "Subject" not in run.text:
                        # Only hyperlink standalone email refs, not "Subject:" lines
                        if len(run.text.strip()) < 50:
                            run.hyperlink.address = "mailto:cobe-info@boisestate.edu?subject=PROGRAM%20REQUEST%20%E2%80%93%20%5BChannel%5D%20%E2%80%93%20%5BDate%20Range%5D"

    # Slide 3 (idx 2) — Golden Rule email
    slide_3 = slides[2]
    for shape in slide_3.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() == "cobe-info@boisestate.edu":
                        run.hyperlink.address = "mailto:cobe-info@boisestate.edu?subject=PROGRAM%20REQUEST%20%E2%80%93%20%5BChannel%5D%20%E2%80%93%20%5BDate%20Range%5D"

    # Slide 45 (idx 44) — Closing email
    slide_45 = slides[44]
    for shape in slide_45.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() == "cobe-info@boisestate.edu":
                        run.hyperlink.address = "mailto:cobe-info@boisestate.edu?subject=PROGRAM%20REQUEST%20%E2%80%93%20%5BChannel%5D%20%E2%80%93%20%5BDate%20Range%5D"

    # Make cobe-info clickable in footers across all slides
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = para.text
                    if "COBE Marketing Playbook" in full_text and "cobe-info" in full_text:
                        # This is a footer — add hyperlink to the email run
                        for run in para.runs:
                            if "cobe-info@boisestate.edu" in run.text and len(run.text) < 80:
                                pass  # Footer is fine as text-only for readability

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 4: NEW SLIDE — Hands-Off Quick Checklist & Lead Times
    # Insert after slide 20 (idx 19, Asset Specs), before Case Study
    # ═══════════════════════════════════════════════════════════════════════
    print("Phase 4: Creating Hands-Off Quick Checklist slide...")

    s_ho_checklist = add_blank_slide(prs)

    # Title
    add_textbox(s_ho_checklist, Inches(0.5), Inches(0.25), Inches(9), Inches(0.55),
                text="Hands-Off: Quick Checklist & Lead Times",
                font_size=32, bold=True, color=DARK_GRAY)

    # Accent line
    add_rectangle(s_ho_checklist, Inches(0.5), Inches(0.78), Inches(1.5), Inches(0.04),
                  fill_color=BSU_BLUE)

    # Subtitle
    add_textbox(s_ho_checklist, Inches(0.5), Inches(0.85), Inches(9), Inches(0.35),
                text="Your one-email summary \u2014 gather these items, then email cobe-info@boisestate.edu",
                font_size=12, color=MED_GRAY)

    # PATH TAG
    add_path_tag(s_ho_checklist, 1, "PATH 1: HANDS-OFF")

    # Left column: Checklist
    cl_box = add_rounded_rect(s_ho_checklist,
                              Inches(0.5), Inches(1.35), Inches(4.25), Inches(3.5),
                              fill_color=LIGHT_GRAY, line_color=BSU_BLUE)
    tf = cl_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "WHAT TO INCLUDE IN YOUR EMAIL"
    set_font(r, size=12, bold=True, color=BSU_BLUE)

    checklist_items = [
        "\u2610  Channel(s) requested (Instagram, LinkedIn, Email, Screens, etc.)",
        "\u2610  Target publish/send date(s) and cadence",
        "\u2610  Event/program details: Who, What, When, Where",
        "\u2610  Photos or images (min 1080px resolution)",
        "\u2610  Registration / RSVP link",
        "\u2610  Campaign objective + success metric",
        "\u2610  Copy draft (\u226420 words) or let COBE write it",
        "\u2610  Budget code (FDCC) if printing is needed",
        "\u2610  Lead owner for approvals",
    ]
    for item in checklist_items:
        add_para(tf, item, font_size=10, color=DARK_GRAY, space_after=Pt(2))

    # Right column: Lead times
    lt_box = add_rounded_rect(s_ho_checklist,
                              Inches(5.1), Inches(1.35), Inches(4.25), Inches(3.5),
                              fill_color=LIGHT_GRAY, line_color=BSU_BLUE)
    tf2 = lt_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Pt(10)
    tf2.margin_right = Pt(10)
    tf2.margin_top = Pt(8)

    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "LEAD TIMES (HANDS-OFF)"
    set_font(r2, size=12, bold=True, color=BSU_BLUE)

    lead_items = [
        ("Instagram / LinkedIn", "10 business days"),
        ("Email Newsletter", "10\u201315 business days"),
        ("Web Updates", "10 business days"),
        ("New Landing Pages", "15 business days"),
        ("Classroom Slides", "5 business days"),
        ("MBEB Screens", "5 business days"),
        ("Campus Screens (other)", "1.5\u20134 weeks"),
        ("Events / Career Fairs", "3\u20136 weeks"),
    ]
    for channel, time in lead_items:
        p_lt = tf2.add_paragraph()
        p_lt.space_before = Pt(1)
        p_lt.space_after = Pt(1)
        r_ch = p_lt.add_run()
        r_ch.text = f"{channel}:  "
        set_font(r_ch, size=10, bold=True, color=DARK_GRAY)
        r_tm = p_lt.add_run()
        r_tm.text = time
        set_font(r_tm, size=10, bold=False, color=DARK_GRAY)

    # Bottom banner
    ban = add_rounded_rect(s_ho_checklist,
                           Inches(0.5), Inches(4.95), Inches(9), Inches(0.25),
                           fill_color=BSU_BLUE)
    tf_ban = ban.text_frame
    tf_ban.margin_top = Pt(2)
    tf_ban.margin_bottom = Pt(2)
    p_ban = tf_ban.paragraphs[0]
    p_ban.alignment = PP_ALIGN.CENTER
    r_ban = p_ban.add_run()
    r_ban.text = "Subject: PROGRAM REQUEST \u2013 [Channel] \u2013 [Date Range]  \u2192  cobe-info@boisestate.edu"
    set_font(r_ban, size=10, bold=True, color=WHITE)
    r_ban.hyperlink.address = "mailto:cobe-info@boisestate.edu?subject=PROGRAM%20REQUEST%20%E2%80%93%20%5BChannel%5D%20%E2%80%93%20%5BDate%20Range%5D"

    # Move this slide to position after slide 20 (0-indexed position 20)
    insert_slide_at(prs, s_ho_checklist, 20)

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 5: NEW SLIDE — HANDS-ON QUICK LINKS
    # Insert after PATH 2 divider (now at idx 24 after Phase 4 insertion)
    # ═══════════════════════════════════════════════════════════════════════
    print("Phase 5: Creating HANDS-ON QUICK LINKS slide...")

    s_quick_links = add_blank_slide(prs)

    # Title
    add_textbox(s_quick_links, Inches(0.5), Inches(0.25), Inches(9), Inches(0.55),
                text="Hands-On Quick Links",
                font_size=32, bold=True, color=DARK_GRAY)

    add_rectangle(s_quick_links, Inches(0.5), Inches(0.78), Inches(1.5), Inches(0.04),
                  fill_color=BSU_ORANGE)

    add_textbox(s_quick_links, Inches(0.5), Inches(0.85), Inches(9), Inches(0.35),
                text="Clickable shortcuts for the Hands-On path \u2014 bookmark this slide",
                font_size=12, color=MED_GRAY)

    add_path_tag(s_quick_links, 2, "PATH 2: HANDS-ON")

    # ── Group 1: START HERE ──
    add_textbox(s_quick_links, Inches(0.5), Inches(1.3), Inches(2.1), Inches(0.3),
                text="START HERE", font_size=11, bold=True, color=BSU_ORANGE)

    make_button_shape(s_quick_links,
                      Inches(0.5), Inches(1.6), Inches(4.25), Inches(0.45),
                      "\u2709  Email Request (pre-filled subject line)",
                      "mailto:cobe-info@boisestate.edu?subject=PROGRAM%20REQUEST%20%E2%80%93%20%5BChannel%5D%20%E2%80%93%20%5BDate%20Range%5D",
                      fill_color=BSU_BLUE, font_size=10)

    # ── Group 2: TEMPLATES ──
    add_textbox(s_quick_links, Inches(0.5), Inches(2.2), Inches(4.25), Inches(0.25),
                text="TEMPLATES", font_size=11, bold=True, color=BSU_ORANGE)

    make_button_shape(s_quick_links,
                      Inches(0.5), Inches(2.5), Inches(4.25), Inches(0.4),
                      "\U0001F4C2  Template Repository (Paste Link)",
                      "https://placeholder.example.com/template-repository",
                      fill_color=BSU_BLUE, font_size=9)

    make_button_shape(s_quick_links,
                      Inches(0.5), Inches(2.95), Inches(4.25), Inches(0.4),
                      "\U0001F4C5  Social Media Calendar Template (Paste Link)",
                      "https://placeholder.example.com/social-calendar-template",
                      fill_color=BSU_BLUE, font_size=9)

    make_button_shape(s_quick_links,
                      Inches(0.5), Inches(3.4), Inches(4.25), Inches(0.4),
                      "\U0001F4CB  Social Media Calendar Example (Paste Link)",
                      "https://placeholder.example.com/social-calendar-example",
                      fill_color=BSU_BLUE, font_size=9)

    make_button_shape(s_quick_links,
                      Inches(0.5), Inches(3.85), Inches(4.25), Inches(0.4),
                      "\U0001F4B0  Budget Workbook Example (Paste Link)",
                      "https://placeholder.example.com/budget-workbook",
                      fill_color=BSU_BLUE, font_size=9)

    # ── Group 3: TRACKING ──
    add_textbox(s_quick_links, Inches(5.25), Inches(1.3), Inches(4.25), Inches(0.3),
                text="TRACKING", font_size=11, bold=True, color=BSU_ORANGE)

    make_button_shape(s_quick_links,
                      Inches(5.25), Inches(1.6), Inches(4.25), Inches(0.45),
                      "\U0001F517  Flowcode QR Generator",
                      "https://www.flowcode.com",
                      fill_color=BSU_ORANGE, font_size=10)

    make_button_shape(s_quick_links,
                      Inches(5.25), Inches(2.15), Inches(4.25), Inches(0.45),
                      "\U0001F4CA  Google Campaign URL Builder (UTMs)",
                      "https://ga-dev-tools.google/campaign-url-builder/",
                      fill_color=BSU_ORANGE, font_size=10)

    # ── Group 4: WEB ──
    add_textbox(s_quick_links, Inches(5.25), Inches(2.85), Inches(4.25), Inches(0.25),
                text="WEB", font_size=11, bold=True, color=BSU_ORANGE)

    make_button_shape(s_quick_links,
                      Inches(5.25), Inches(3.15), Inches(4.25), Inches(0.45),
                      "\U0001F393  WordPress Accessibility Training",
                      "https://www.boisestate.edu/webguide/wordpress-accessibility-compliance-training/",
                      fill_color=RGBColor(0x44, 0x44, 0x44), font_size=10)

    make_button_shape(s_quick_links,
                      Inches(5.25), Inches(3.7), Inches(4.25), Inches(0.45),
                      "\U0001F4E7  Helpdesk Support",
                      "mailto:helpdesk@boisestate.edu",
                      fill_color=RGBColor(0x44, 0x44, 0x44), font_size=10)

    # Placeholder note
    note_box = add_textbox(s_quick_links,
                           Inches(0.5), Inches(4.5), Inches(9), Inches(0.4),
                           text="Links marked \u201c(Paste Link)\u201d are placeholders \u2014 replace with your live URLs before distributing.",
                           font_size=9, color=MED_GRAY, alignment=PP_ALIGN.LEFT)

    # Insert after PATH 2 divider. After Phase 4, PATH 2 divider is at idx 24.
    # We want this new slide at position 25 (0-indexed).
    insert_slide_at(prs, s_quick_links, 25)

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 6: NEW SLIDE — Asset Pack Folder Structure
    # Insert after Hands-On Submission Checklist (was idx 34, now idx 36 after 2 inserts)
    # ═══════════════════════════════════════════════════════════════════════
    print("Phase 6: Creating Asset Pack Folder Structure slide...")

    s_folder = add_blank_slide(prs)

    add_textbox(s_folder, Inches(0.5), Inches(0.25), Inches(9), Inches(0.55),
                text="Hands-On Asset Pack Folder Structure",
                font_size=32, bold=True, color=DARK_GRAY)

    add_rectangle(s_folder, Inches(0.5), Inches(0.78), Inches(1.5), Inches(0.04),
                  fill_color=BSU_ORANGE)

    add_textbox(s_folder, Inches(0.5), Inches(0.85), Inches(9), Inches(0.35),
                text="Organize your deliverables before submitting to COBE for review",
                font_size=12, color=MED_GRAY)

    add_path_tag(s_folder, 2, "PATH 2: HANDS-ON")

    # Left: Folder tree
    folder_box = add_rounded_rect(s_folder,
                                  Inches(0.5), Inches(1.35), Inches(4.5), Inches(3.3),
                                  fill_color=LIGHT_GRAY, line_color=BSU_ORANGE)
    tf_f = folder_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = Pt(12)
    tf_f.margin_right = Pt(10)
    tf_f.margin_top = Pt(8)

    p_f = tf_f.paragraphs[0]
    r_f = p_f.add_run()
    r_f.text = "RECOMMENDED FOLDER STRUCTURE"
    set_font(r_f, size=12, bold=True, color=BSU_ORANGE)

    folders = [
        ("\U0001F4C1  /01_Copy", "Captions, email body text, talking points"),
        ("\U0001F4C1  /02_Graphics", "Final images at spec (1080\u00D71080, 1080\u00D71920, 1920\u00D71080)"),
        ("\U0001F4C1  /03_Links_Tracking", "UTM URLs, QR code files, Flowcode exports"),
        ("\U0001F4C1  /04_Schedule", "Completed social media calendar (.xlsx)"),
        ("\U0001F4C1  /05_Budget (optional)", "Budget proposal / workbook if using paid channels"),
    ]
    for folder_name, desc in folders:
        p_fn = tf_f.add_paragraph()
        p_fn.space_before = Pt(6)
        p_fn.space_after = Pt(0)
        r_fn = p_fn.add_run()
        r_fn.text = folder_name
        set_font(r_fn, size=11, bold=True, color=DARK_GRAY)

        p_desc = tf_f.add_paragraph()
        p_desc.space_before = Pt(0)
        p_desc.space_after = Pt(1)
        r_desc = p_desc.add_run()
        r_desc.text = f"     {desc}"
        set_font(r_desc, size=9, color=MED_GRAY)

    # Right: Naming conventions
    name_box = add_rounded_rect(s_folder,
                                Inches(5.25), Inches(1.35), Inches(4.25), Inches(3.3),
                                fill_color=LIGHT_GRAY, line_color=BSU_ORANGE)
    tf_n = name_box.text_frame
    tf_n.word_wrap = True
    tf_n.margin_left = Pt(12)
    tf_n.margin_right = Pt(10)
    tf_n.margin_top = Pt(8)

    p_n = tf_n.paragraphs[0]
    r_n = p_n.add_run()
    r_n.text = "FILE NAMING CONVENTION"
    set_font(r_n, size=12, bold=True, color=BSU_ORANGE)

    add_para(tf_n, "Pattern:", font_size=10, bold=True, color=DARK_GRAY,
             space_before=Pt(6))
    add_para(tf_n, "topic_start-end-date_size.ext", font_size=11, bold=True,
             color=BSU_BLUE)

    add_para(tf_n, "", font_size=4)
    add_para(tf_n, "Examples:", font_size=10, bold=True, color=DARK_GRAY,
             space_before=Pt(4))
    examples = [
        "CareerFair_2025-10-01_1920x1080.jpg",
        "RHM-ShuttleStop-Feb25_QR.png",
        "RHM_Newsletter_Feb15-28_Copy.docx",
        "RHM_SocialCalendar_Spring2025.xlsx",
    ]
    for ex in examples:
        add_para(tf_n, f"\u2022  {ex}", font_size=9, color=DARK_GRAY)

    add_para(tf_n, "", font_size=4)
    add_para(tf_n, "Tips:", font_size=10, bold=True, color=DARK_GRAY,
             space_before=Pt(4))
    tips = [
        "No spaces in file names \u2014 use hyphens or underscores",
        "Include dimensions in graphic file names",
        "Use date format YYYY-MM-DD for sorting",
    ]
    for tip in tips:
        add_para(tf_n, f"\u2022  {tip}", font_size=9, color=DARK_GRAY)

    # Bottom banner
    ban2 = add_rounded_rect(s_folder,
                            Inches(0.5), Inches(4.75), Inches(9), Inches(0.45),
                            fill_color=BSU_ORANGE)
    tf_ban2 = ban2.text_frame
    tf_ban2.margin_top = Pt(4)
    tf_ban2.margin_bottom = Pt(4)
    tf_ban2.margin_left = Pt(10)
    p_ban2 = tf_ban2.paragraphs[0]
    p_ban2.alignment = PP_ALIGN.CENTER
    r_ban2 = p_ban2.add_run()
    r_ban2.text = "Zip your folder and attach it to your email to cobe-info@boisestate.edu"
    set_font(r_ban2, size=11, bold=True, color=WHITE)

    # Insert after Hands-On Submission Checklist.
    # Original slide 35 (idx 34). After 2 prior inserts, it's now at idx 36.
    # We want this new slide at position 37.
    insert_slide_at(prs, s_folder, 37)

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 7: Update slide numbers in footers
    # ═══════════════════════════════════════════════════════════════════════
    print("Phase 7: Updating slide numbers...")

    all_slides = list(prs.slides)
    for idx, slide in enumerate(all_slides):
        slide_num = idx + 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                # Detect page number shapes: small text box with just a number
                # positioned near bottom-right
                if (txt.isdigit() and
                    shape.top and shape.top > Inches(4.5) and
                    shape.left and shape.left > Inches(7)):
                    # Update the number
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip().isdigit():
                                run.text = str(slide_num)

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 8: Add speaker notes where slides are dense
    # ═══════════════════════════════════════════════════════════════════════
    print("Phase 8: Adding speaker notes to dense slides...")

    notes_map = {
        1: "Welcome to the COBE Marketing Playbook. This guide covers both Hands-Off and Hands-On paths for marketing through COBE. All requests start at cobe-info@boisestate.edu.",
        4: "Two paths: Hands-Off (concierge \u2014 send facts, COBE handles everything) and Hands-On (DIY with templates, COBE reviews and approves). Both start with an email to cobe-info@boisestate.edu.",
        12: "This is the core Hands-Off workflow. Emphasize that Step 2 (the email) is all the requester needs to do. COBE handles steps 3 and 4.",
        13: "Walk through the required vs. helpful additions. Stress that a complete first email avoids back-and-forth and speeds turnaround.",
        25: "This is the core Hands-On workflow. Emphasize that steps 1-3 are done by the requester using COBE templates, then step 4 is the COBE review gate.",
    }

    # Adjust indices for new slides inserted (slides 1-20 unchanged, 21 is new,
    # 22-25 shifted +1, 26 is new, 27-37 shifted +2, 38 is new, 39+ shifted +3)
    # Map from desired slide number to actual 0-indexed position
    for slide_num_1based, note_text in notes_map.items():
        # For the original slide numbers, map to new positions
        idx = slide_num_1based - 1  # 0-indexed
        if idx < len(all_slides):
            slide = all_slides[idx]
            if not slide.has_notes_slide:
                slide.notes_slide  # Creates notes slide
            notes_tf = slide.notes_slide.notes_text_frame
            if not notes_tf.text.strip():
                notes_tf.paragraphs[0].text = note_text

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 9: Expand Digital Screens & Physical Materials guidance
    # Add mini-checklist as speaker notes to slides 33 and 34 (now shifted)
    # ═══════════════════════════════════════════════════════════════════════
    print("Phase 9: Enhancing Digital Screens & Physical Materials content...")

    # After 3 inserts, original slide 33 (idx 32) is now at a new position.
    # Original slides 25-35 shifted by +2 (two inserts before them: at positions 20 and 25)
    # So original slide 33 (idx 32) -> now idx 34, original 34 -> idx 35

    # Find the slides by title text to be safe
    for idx, slide in enumerate(all_slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Building Flyers, Posters & Print Materials" in shape.text_frame.text:
                    # Add speaker notes with QR placement guidance
                    if not slide.notes_slide.notes_text_frame.text.strip():
                        slide.notes_slide.notes_text_frame.paragraphs[0].text = (
                            "Physical Materials Quick Reference:\n"
                            "- Best for: Event promotion, tabling, classroom visits\n"
                            "- QR code placement: Lower-right corner, minimum 1\u201d\u00D71\u201d for scannability\n"
                            "- Confirm all print specs with COBE Marketing before ordering\n"
                            "- MBEB bulletin boards only for flyers (not walls/windows/doors)\n"
                            "- Always include a unique QR code per placement for tracking\n"
                            "- Return poster boards and easels to the Dean's Office after use"
                        )
                    break

    for idx, slide in enumerate(all_slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Campus Digital Screens: Complete Guide" in shape.text_frame.text:
                    if not slide.notes_slide.notes_text_frame.text.strip():
                        slide.notes_slide.notes_text_frame.paragraphs[0].text = (
                            "Digital Screens Quick Reference:\n"
                            "- Best for: Brand awareness, event countdowns, daily visibility\n"
                            "- MBEB screens: 1920\u00D71080 JPG, dark background, NO logo overlay (Clean Screen Rule)\n"
                            "- Submit 3\u20134 weeks before desired start date\n"
                            "- Name files: topic_start-end-date_size.jpg\n"
                            "- Confirm specs with COBE Marketing for any non-MBEB campus screens\n"
                            "- Each campus location has different contacts and dimensions \u2014 see this slide for the full list"
                        )
                    break

    # ═══════════════════════════════════════════════════════════════════════
    # PHASE 10: Save
    # ═══════════════════════════════════════════════════════════════════════
    print(f"Saving to {OUTPUT_FILE}...")
    prs.save(OUTPUT_FILE)

    final_count = len(list(prs.slides))
    print(f"Done! {final_count} slides saved ({final_count - total_original} new slides added).")
    print(f"Output: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
