#!/usr/bin/env python3
"""
Comprehensive update script for COBE Marketing Playbook.
Transforms "COBE Marketing Playbook Updated (1).pptx" into version (2).
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ─── Constants ────────────────────────────────────────────────────────────────

SRC = "COBE Marketing Playbook Updated (1).pptx"
DST = "COBE Marketing Playbook Updated (2).pptx"

# Brand colors
BOISE_BLUE = RGBColor(0x00, 0x33, 0xA0)
BOISE_ORANGE = RGBColor(0xD6, 0x43, 0x09)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
TAG_GREEN = RGBColor(0x00, 0x80, 0x60)
TAG_BLUE = RGBColor(0x00, 0x55, 0xA0)
TAG_ORANGE = RGBColor(0xD6, 0x43, 0x09)

# Slide dimensions in EMU
SLIDE_W = 9144000
SLIDE_H = 5143500

# ─── Helper functions ─────────────────────────────────────────────────────────

def find_shape(slide, name):
    """Find a shape by its exact name."""
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def find_shape_containing(slide, text_fragment):
    """Find first shape containing a text fragment."""
    for s in slide.shapes:
        if s.has_text_frame:
            full = s.text_frame.text
            if text_fragment in full:
                return s
    return None

def set_text_preserve_first_run(shape, new_text):
    """Set shape text, preserving formatting from the first run of the first paragraph."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        font_name = first_run.font.name
        font_size = first_run.font.size
        font_bold = first_run.font.bold
        font_color = first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else None
        # Clear all paragraphs and set new text
        for p in tf.paragraphs:
            p.clear()
        tf.paragraphs[0].text = new_text
        if tf.paragraphs[0].runs:
            r = tf.paragraphs[0].runs[0]
            if font_name: r.font.name = font_name
            if font_size: r.font.size = font_size
            if font_bold is not None: r.font.bold = font_bold
            if font_color: r.font.color.rgb = font_color
    else:
        tf.text = new_text

def set_para_text(shape, para_idx, new_text):
    """Set text of a specific paragraph in a shape, preserving run formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if para_idx >= len(tf.paragraphs):
        return
    p = tf.paragraphs[para_idx]
    if p.runs:
        # Preserve formatting from first run
        first_run = p.runs[0]
        font_name = first_run.font.name
        font_size = first_run.font.size
        font_bold = first_run.font.bold
        font_color = first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else None
        p.clear()
        run = p.add_run()
        run.text = new_text
        if font_name: run.font.name = font_name
        if font_size: run.font.size = font_size
        if font_bold is not None: run.font.bold = font_bold
        if font_color: run.font.color.rgb = font_color
    else:
        p.text = new_text

def replace_in_all_runs(shape, old, new):
    """Replace text in all runs of all paragraphs of a shape."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                changed = True
    return changed

def deck_wide_replace(prs, old, new):
    """Replace text across all slides."""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if replace_in_all_runs(shape, old, new):
                    count += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)
                                        count += 1
    return count

def add_textbox(slide, left, top, width, height, text, font_size=10,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name='Calibri', bg_color=None):
    """Add a text box to a slide with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    if bg_color:
        from pptx.oxml.ns import qn
        sp = txBox._element
        spPr = sp.find(qn('p:txBody'))
        # Use shape fill instead
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox

def add_rounded_box(slide, left, top, width, height, text_lines,
                    header_text=None, font_size=8, bg_color=None,
                    header_color=WHITE, text_color=BLACK, header_bg=None):
    """Add a rounded rectangle with optional header and body text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color or LIGHT_GRAY
    shape.line.color.rgb = MED_GRAY
    shape.line.width = Pt(0.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)

    if header_text:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = header_text
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = header_color if header_bg else text_color
        run.font.name = 'Calibri'
        p.space_after = Pt(2)

    for i, line in enumerate(text_lines):
        if header_text or i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.name = 'Calibri'
        p.space_after = Pt(1)

    return shape

def add_funnel_tag(slide, left, top, width, height, text, bg_color):
    """Add a small funnel indicator tag."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(7)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'
    return shape

def add_activation_boxes(slide, left, top, width, channel_items_ho, channel_items_hn):
    """Add Hands-Off and Hands-On activation boxes to a channel slide."""
    box_height = Emu(850000)
    gap = Emu(60000)

    # Header
    add_textbox(slide, left, top, width, Emu(180000),
                "HOW TO ACTIVATE", font_size=9, bold=True, color=DARK_GRAY)

    # Hands-Off box
    ho_top = top + Emu(190000)
    ho_lines = [
        "Email COBE Marketing at cobe-info@boisestate.edu with:"
    ] + channel_items_ho
    add_rounded_box(slide, left, ho_top, width, box_height,
                   ho_lines, header_text="Hands-Off (Concierge)",
                   font_size=7, bg_color=RGBColor(0xE8, 0xEE, 0xF7),
                   text_color=DARK_GRAY)

    # Hands-On box
    hn_top = ho_top + box_height + gap
    hn_lines = [
        "Create using COBE templates (see Templates & Links Hub).",
        "Email cobe-info@boisestate.edu for review/approval with:"
    ] + channel_items_hn
    add_rounded_box(slide, left, hn_top, width, box_height,
                   hn_lines, header_text="Hands-On (DIY w/ Support)",
                   font_size=7, bg_color=RGBColor(0xFD, 0xF0, 0xE5),
                   text_color=DARK_GRAY)

def add_screenshot_placeholder(slide, left, top, width, height, label):
    """Add a placeholder box where a screenshot should go."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.color.rgb = MED_GRAY
    shape.line.width = Pt(1)
    shape.line.dash_style = 2  # dash

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(8)
    run.font.color.rgb = MED_GRAY
    run.font.name = 'Calibri'
    run.font.italic = True
    return shape

def delete_slide(prs, index):
    """Delete a slide by index."""
    rId = prs.slides._sldIdLst[index].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        # Try the 'r:id' attribute
        rId = prs.slides._sldIdLst[index].attrib.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    from lxml import etree
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    sldIdLst.remove(sldId)

def add_multi_line_textbox(slide, left, top, width, height, lines_config):
    """Add a textbox with multiple paragraphs, each with its own formatting.
    lines_config is a list of dicts: {text, size, bold, color, alignment}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, cfg in enumerate(lines_config):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = cfg.get('alignment', PP_ALIGN.LEFT)
        p.space_after = Pt(cfg.get('space_after', 2))
        run = p.add_run()
        run.text = cfg['text']
        run.font.size = Pt(cfg.get('size', 10))
        run.font.bold = cfg.get('bold', False)
        run.font.color.rgb = cfg.get('color', BLACK)
        run.font.name = cfg.get('font', 'Calibri')
    return txBox

# ─── Main update logic ───────────────────────────────────────────────────────

def main():
    prs = Presentation(SRC)
    slides = list(prs.slides)

    print("Starting deck update...")
    print(f"Total slides: {len(slides)}")

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 1: Deck-wide text replacements
    # ═══════════════════════════════════════════════════════════════════════════
    print("\n--- Phase 1: Deck-wide replacements ---")

    # Replace "us/we" patterns with "COBE Marketing"
    replacements = [
        ("Send us the facts.", "Email COBE Marketing the facts."),
        ("We'll handle the rest.", "COBE Marketing handles copy, design, scheduling, approvals."),
        ("Use our templates.", "Use COBE templates (see Templates & Links Hub)."),
        ("We review & approve.", "COBE Marketing reviews and approves."),
        ("we create them", "COBE Marketing creates them"),
        ("use our templates", "use COBE templates from the Templates & Links Hub"),
        ("Send facts - We create", "Send facts → COBE Marketing creates"),
        ("Use templates - We approve", "Use COBE templates → COBE Marketing approves"),
        ("We can always scale down, but we can't scale up.",
         "COBE Marketing can always scale down but cannot scale up."),
        ("We'll do our best to accommodate.",
         "COBE Marketing will do its best to accommodate."),
        ("How we use different channels",
         "How COBE Marketing uses different channels"),
        ("let's build something great together.",
         "COBE Marketing is here to help you succeed."),
        ("xxx@email", "cobe-info@boisestate.edu"),
    ]

    for old, new in replacements:
        c = deck_wide_replace(prs, old, new)
        if c > 0:
            print(f"  Replaced '{old[:40]}...' → {c} occurrences")

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 2: Slide-by-slide edits
    # ═══════════════════════════════════════════════════════════════════════════
    print("\n--- Phase 2: Slide-by-slide edits ---")

    # ─── Slide 1 (index 0) — Title: minor polish ─────────────────────────────
    print("  Slide 1: Title polish")
    s1 = slides[0]
    # Ensure consistent capitalization on subtitle
    sub_shape = find_shape(s1, 'Google Shape;18;p1')
    if sub_shape:
        set_para_text(sub_shape, 0, "Your Action-Ready Guide to Brand-Consistent Marketing")
        set_para_text(sub_shape, 1, "For Events, Programs, and Partnerships")

    # ─── Slide 2 (index 1) — The Bandwidth Problem ───────────────────────────
    print("  Slide 2: The Bandwidth Problem")
    s2 = slides[1]

    # Replace the challenge quote
    quote_shape = find_shape(s2, 'Google Shape;31;p2')
    if quote_shape:
        set_text_preserve_first_run(quote_shape,
            "I have an event, newsletter item, or announcement to promote\u2014what\u2019s the fastest way to get it out?")

    # Update the explanation paragraph
    explain_shape = find_shape(s2, 'Google Shape;32;p2')
    if explain_shape and explain_shape.has_text_frame:
        set_para_text(explain_shape, 0,
            "Faculty and staff are experts in their fields, not marketing. Learning brand guidelines, specs, and approval workflows takes time you don\u2019t have.")
        if len(explain_shape.text_frame.paragraphs) > 2:
            set_para_text(explain_shape, 2,
                "Student organizations face similar constraints with limited resources.")

    # Reframe COBE MARKETING REACH header
    reach_header = find_shape(s2, 'Google Shape;36;p2')
    if reach_header:
        set_text_preserve_first_run(reach_header,
            "COBE MARKETING FREE CHANNELS")

    # Add clarifying text below reach header - modify the follower counts area
    # The shape at 5212080,1508760 is the "1,436" followers text
    # We want to add context. Let's modify the area above the pathway boxes.
    # Actually - let me add a new textbox below the reach grid.
    # The reach grid ends around y=2700. Pathway boxes start at y=2926.
    # There's space for a small note.
    # No, actually the grid fills the space. Let me just update the subtitle.
    solution_sub = find_shape(s2, 'Google Shape;34;p2')
    if solution_sub:
        set_text_preserve_first_run(solution_sub,
            "COBE Marketing manages free channels reaching thousands of students, alumni, and employers.")

    # Update Hands-Off box text
    ho_title = find_shape(s2, 'Google Shape;46;p2')  # "HANDS-OFF"
    ho_body = find_shape(s2, 'Google Shape;47;p2')   # "Send us the facts..."
    if ho_body and ho_body.has_text_frame:
        set_para_text(ho_body, 0, "Email COBE Marketing the facts.")
        if len(ho_body.text_frame.paragraphs) > 1:
            set_para_text(ho_body, 1, "COBE Marketing handles copy, design, scheduling, approvals.")

    # Update Hands-On box text
    hn_body = find_shape(s2, 'Google Shape;51;p2')  # "Use our templates..."
    if hn_body and hn_body.has_text_frame:
        set_para_text(hn_body, 0, "Use COBE templates to draft assets.")
        if len(hn_body.text_frame.paragraphs) > 1:
            set_para_text(hn_body, 1, "COBE Marketing reviews and approves.")

    # ─── Slide 3 (index 2) — Choose Your Path ────────────────────────────────
    print("  Slide 3: Choose Your Path")
    s3 = slides[2]

    # Update Hands-Off description
    ho_desc = find_shape(s3, 'Google Shape;72;p3')
    if ho_desc:
        set_text_preserve_first_run(ho_desc,
            "For busy faculty who want to send basic facts and have COBE Marketing handle everything else.")

    # "WE HANDLE:" → "COBE MARKETING HANDLES:"
    we_handle = find_shape(s3, 'Google Shape;75;p3')
    if we_handle:
        set_text_preserve_first_run(we_handle, "COBE MARKETING HANDLES:")

    # Update Hands-On description
    hn_desc = find_shape(s3, 'Google Shape;80;p3')
    if hn_desc:
        set_text_preserve_first_run(hn_desc,
            "For users who want creative control and will draft their own copy and graphics using COBE templates.")

    # "WE PROVIDE:" → "COBE MARKETING PROVIDES:"
    we_provide = find_shape(s3, 'Google Shape;83;p3')
    if we_provide:
        set_text_preserve_first_run(we_provide, "COBE MARKETING PROVIDES:")

    # Update what COBE Marketing provides
    provides_detail = find_shape(s3, 'Google Shape;84;p3')
    if provides_detail:
        set_text_preserve_first_run(provides_detail,
            "Templates, Brand Review, Final Approval, Scheduling as Needed")

    # ─── Slide 4 (index 3) — Golden Rule ─────────────────────────────────────
    print("  Slide 4: Golden Rule")
    s4 = slides[3]

    # Update the instruction text
    instr = find_shape(s4, 'Google Shape;94;p4')
    if instr:
        set_text_preserve_first_run(instr,
            "Whether going Hands-Off or Hands-On, every marketing request begins by emailing COBE Marketing at this address.")

    # ─── Slide 5 (index 4) — Quick-Start ─────────────────────────────────────
    print("  Slide 5: Quick-Start")
    s5 = slides[4]

    # Add a line reminding users to match goal → funnel stage → channel
    # The quick tips area starts at y=3657600. Let me add before that.
    # Actually, better to add to the existing QUICK TIPS
    qt_shape = find_shape(s5, 'Google Shape;128;p5')
    if qt_shape and qt_shape.has_text_frame:
        tf = qt_shape.text_frame
        # Add a new paragraph at the end
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Match your goal \u2192 funnel stage \u2192 channel for best results (see Marketing Funnel slide)"
        # Copy formatting from first paragraph's first run
        if tf.paragraphs[0].runs:
            ref = tf.paragraphs[0].runs[0]
            if ref.font.size: run.font.size = ref.font.size
            if ref.font.name: run.font.name = ref.font.name
            if ref.font.color and ref.font.color.rgb: run.font.color.rgb = ref.font.color.rgb

    # ─── Slide 10 (index 9) — Marketing Funnel ───────────────────────────────
    print("  Slide 10: Marketing Funnel")
    s10 = slides[9]

    # Update subtitle to remove "we"
    s10_sub = find_shape(s10, 'Google Shape;200;p10')
    if s10_sub:
        set_text_preserve_first_run(s10_sub,
            "Match your goal to the right funnel stage and channel")

    # Rename "CHANNELS" to "RECOMMENDED CHANNELS" in all three funnel columns
    for name in ['Google Shape;206;p10', 'Google Shape;213;p10', 'Google Shape;220;p10']:
        ch_shape = find_shape(s10, name)
        if ch_shape:
            set_text_preserve_first_run(ch_shape, "RECOMMENDED CHANNELS")

    # Add "How to use this" callout at the bottom
    add_multi_line_textbox(s10, Emu(457200), Emu(4572000), Emu(8229600), Emu(500000), [
        {"text": "HOW TO USE THIS:  ", "size": 8, "bold": True, "color": DARK_GRAY, "space_after": 1},
        {"text": "1) Define your goal (awareness / signups / attendance)   2) Choose the funnel stage   3) Use the recommended channels \u2192 go to each channel page for activation details + specs",
         "size": 7, "color": DARK_GRAY, "space_after": 0},
    ])

    # ─── Slide 11 (index 10) — Instagram (rebuild with channel template) ─────
    print("  Slide 11: Instagram - adding funnel tags, activation, screenshot")
    s11 = slides[10]

    # Add funnel tags top-right
    tag_y = Emu(274320)
    tag_w = Emu(750000)
    tag_h = Emu(220000)
    add_funnel_tag(s11, Emu(5943600), tag_y, tag_w, tag_h, "Awareness", TAG_GREEN)
    add_funnel_tag(s11, Emu(5943600) + tag_w + Emu(60000), tag_y, tag_w + Emu(150000), tag_h, "Consideration", TAG_BLUE)

    # Shrink/reposition the SPECS box to make room
    specs_shape = find_shape(s11, 'Google Shape;243;p11')
    if specs_shape:
        specs_shape.top = Emu(1645920)
        specs_shape.height = Emu(880000)
    specs_header = find_shape(s11, 'Google Shape;244;p11')
    if specs_header:
        specs_header.top = Emu(1700000)
    specs_body = find_shape(s11, 'Google Shape;245;p11')
    if specs_body:
        specs_body.top = Emu(1870000)
        specs_body.height = Emu(600000)

    # Shrink PRO TIPS box
    pro_shape = find_shape(s11, 'Google Shape;246;p11')
    if pro_shape:
        pro_shape.top = Emu(2580000)
        pro_shape.height = Emu(600000)
    pro_header = find_shape(s11, 'Google Shape;247;p11')
    if pro_header:
        pro_header.top = Emu(2610000)
    pro_body = find_shape(s11, 'Google Shape;248;p11')
    if pro_body:
        pro_body.top = Emu(2780000)
        pro_body.height = Emu(380000)
        # Add lead time to pro tips
        tf = pro_body.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Lead time: 10 business days"
        if tf.paragraphs[0].runs:
            ref = tf.paragraphs[0].runs[0]
            if ref.font.size: run.font.size = ref.font.size
            if ref.font.name: run.font.name = ref.font.name
            if ref.font.color and ref.font.color.rgb: run.font.color.rgb = ref.font.color.rgb

    # Shrink the bottom tip bar to make room for activation
    tip_bar_bg = find_shape(s11, 'Google Shape;249;p11')
    if tip_bar_bg:
        tip_bar_bg.top = Emu(4750000)
        tip_bar_bg.height = Emu(350000)
    tip_bar_text = find_shape(s11, 'Google Shape;250;p11')
    if tip_bar_text:
        tip_bar_text.top = Emu(4750000)
        tip_bar_text.height = Emu(350000)

    # Add activation boxes in the space we freed (right side, below pro tips)
    add_activation_boxes(s11,
        left=Emu(5029200), top=Emu(3250000), width=Emu(3657600),
        channel_items_ho=[
            "\u2022 Event/program details (who, what, when, where)",
            "\u2022 Photos (min 1080px)",
            "\u2022 Registration/RSVP link",
        ],
        channel_items_hn=[
            "\u2022 Finished graphic (1080x1080 or 1080x1920)",
            "\u2022 Caption text + CTA",
        ]
    )

    # Add screenshot placeholder (between content types and activation)
    # Place it in the left column area below the Grid Posts content
    add_screenshot_placeholder(s11,
        Emu(457200), Emu(3800000), Emu(4400000), Emu(880000),
        "[Placement Screenshot: Instagram feed/story example]"
    )

    # ─── Slide 12 (index 11) — Instagram duplicate: fix for deletion ─────────
    # We'll delete this slide after all edits (in Phase 3)
    print("  Slide 12: Marked for deletion")

    # ─── Slide 13 (index 12) — LinkedIn (rebuild with channel template) ──────
    print("  Slide 13: LinkedIn - adding funnel tags, activation, screenshot")
    s13 = slides[12]

    # Add funnel tags
    add_funnel_tag(s13, Emu(5943600), tag_y, tag_w + Emu(150000), tag_h, "Consideration", TAG_BLUE)
    add_funnel_tag(s13, Emu(5943600) + tag_w + Emu(210000), tag_y, tag_w, tag_h, "Conversion", TAG_ORANGE)

    # Shrink KEY STRATEGY box to make room
    ks_bg = find_shape(s13, 'Google Shape;297;p12')
    if ks_bg:
        ks_bg.height = Emu(700000)
    ks_header = find_shape(s13, 'Google Shape;298;p12')
    ks_body1 = find_shape(s13, 'Google Shape;299;p12')
    ks_body2 = find_shape(s13, 'Google Shape;300;p12')

    # Shrink SPECS box
    sp_bg = find_shape(s13, 'Google Shape;301;p12')
    if sp_bg:
        sp_bg.top = Emu(2420000)
        sp_bg.height = Emu(700000)
    sp_header = find_shape(s13, 'Google Shape;302;p12')
    if sp_header:
        sp_header.top = Emu(2450000)
    sp_body = find_shape(s13, 'Google Shape;303;p12')
    if sp_body:
        sp_body.top = Emu(2600000)
        sp_body.height = Emu(480000)
        # Add lead time info
        tf = sp_body.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Lead time: 10 business days (events: 3 weeks)"
        if tf.paragraphs[0].runs:
            ref = tf.paragraphs[0].runs[0]
            if ref.font.size: run.font.size = ref.font.size
            if ref.font.name: run.font.name = ref.font.name
            if ref.font.color and ref.font.color.rgb: run.font.color.rgb = ref.font.color.rgb

    # Shrink Faculty Advocacy bar
    fa_bg = find_shape(s13, 'Google Shape;304;p12')
    if fa_bg:
        fa_bg.top = Emu(4750000)
        fa_bg.height = Emu(350000)
    fa_text = find_shape(s13, 'Google Shape;305;p12')
    if fa_text:
        fa_text.top = Emu(4750000)
        fa_text.height = Emu(350000)

    # Add activation boxes
    add_activation_boxes(s13,
        left=Emu(5029200), top=Emu(3200000), width=Emu(3657600),
        channel_items_ho=[
            "\u2022 Story/achievement details + outcome data",
            "\u2022 Partner names to tag",
            "\u2022 Photos/headshots + registration link",
        ],
        channel_items_hn=[
            "\u2022 Finished post copy + graphics (1920x1080)",
            "\u2022 Partner tags + links",
        ]
    )

    # Add screenshot placeholder
    add_screenshot_placeholder(s13,
        Emu(457200), Emu(3800000), Emu(4400000), Emu(880000),
        "[Placement Screenshot: LinkedIn post/article example]"
    )

    # ─── Slide 14 (index 13) — Email Marketing (rebuild as channel template) ─
    print("  Slide 14: Email Marketing - restructuring as channel template")
    s14 = slides[13]

    # Add funnel tags
    add_funnel_tag(s14, Emu(7300000), Emu(274320), tag_w, tag_h, "Conversion", TAG_ORANGE)

    # Restructure: convert the essay-like "Power of Email" and "Burnout Risk"
    # into more scannable format.
    # Power of Email box - make it briefer
    power_body = find_shape(s14, 'Google Shape;317;p13')
    if power_body:
        set_text_preserve_first_run(power_body,
            "Email has the highest conversion rate of all COBE channels. Direct inbox access means higher engagement for RSVPs, applications, and deadlines.")

    burnout_body = find_shape(s14, 'Google Shape;320;p13')
    if burnout_body:
        set_text_preserve_first_run(burnout_body,
            "Over-emailing causes unsubscribes. Only send when you have genuine value. Coordinate with COBE Marketing on frequency.")

    # Rename EMAIL RULES to PRO TIPS
    rules_header = find_shape(s14, 'Google Shape;321;p13')
    if rules_header:
        set_text_preserve_first_run(rules_header, "PRO TIPS")

    # Update newsletter strategy box to include lead time
    newsletter_header = find_shape(s14, 'Google Shape;329;p13')
    if newsletter_header:
        set_text_preserve_first_run(newsletter_header,
            "SPECS & LEAD TIME")
    newsletter_body = find_shape(s14, 'Google Shape;330;p13')
    if newsletter_body:
        set_text_preserve_first_run(newsletter_body,
            "Lead time: 10\u201315 business days. Copy limit ~100 words per item. Include: final link/CTA, copy, images with alt text. Subscribers: 5,000+.")

    # Add activation boxes (overlay on the left-bottom area)
    add_activation_boxes(s14,
        left=Emu(457200), top=Emu(3100000), width=Emu(3800000),
        channel_items_ho=[
            "\u2022 Event/program details",
            "\u2022 Copy (\u2264100 words), CTA, final link",
            "\u2022 Images with alt text",
        ],
        channel_items_hn=[
            "\u2022 Full copy draft + images",
            "\u2022 Target segment (students/alumni/employers)",
        ]
    )

    # Shift Email Rules items up to make room
    # These are at y=3383280. Move them down as they'll be part of pro tips in the right area

    # Add screenshot placeholder
    add_screenshot_placeholder(s14,
        Emu(4754880), Emu(3108960), Emu(3931920), Emu(780000),
        "[Placement Screenshot: Newsletter section example]"
    )

    # ─── Slide 15 (index 14) — Internal Channels (rebuild with template) ─────
    print("  Slide 15: Internal Channels - adding funnel tags, activation, screenshot")
    s15 = slides[14]

    # Add funnel tags
    add_funnel_tag(s15, Emu(6300000), Emu(274320), tag_w, tag_h, "Awareness", TAG_GREEN)
    add_funnel_tag(s15, Emu(6300000) + tag_w + Emu(60000), Emu(274320), tag_w, tag_h, "Conversion", TAG_ORANGE)

    # Add activation boxes below the main content
    add_activation_boxes(s15,
        left=Emu(4754880), top=Emu(3200000), width=Emu(3931920),
        channel_items_ho=[
            "\u2022 Event details, dates, QR code destination URL",
            "\u2022 Photos/graphics (if available)",
        ],
        channel_items_hn=[
            "\u2022 Finished slide (1920x1080, dark bg, no logo)",
            "\u2022 QR code linking to registration/RSVP",
        ]
    )

    # Add lead time + specs note
    add_multi_line_textbox(s15, Emu(457200), Emu(4600000), Emu(8229600), Emu(450000), [
        {"text": "LEAD TIME: Classroom slides: 5 business days  |  MBEB screens: 5 business days  |  SPECS: 1920x1080px, dark background, no logo overlays, large QR codes",
         "size": 7, "bold": False, "color": DARK_GRAY},
    ])

    # Add screenshot placeholder (in the space to the right of the campus screens table)
    add_screenshot_placeholder(s15,
        Emu(4754880), Emu(3200000) - Emu(400000), Emu(3931920), Emu(400000),
        "[Placement Screenshot: Classroom slide / MBEB signage example]"
    )

    # ─── Slide 19 (index 18) — Sample Request Email ──────────────────────────
    print("  Slide 19: Sample Request Email")
    s19 = slides[18]

    # Update title to specify this is a Hands-Off sample
    title_shape = find_shape(s19, 'Google Shape;418;p18')
    if title_shape:
        set_text_preserve_first_run(title_shape, "Sample Request Email (Hands-Off)")

    # Update subtitle
    sub_shape = find_shape(s19, 'Google Shape;420;p18')
    if sub_shape:
        set_text_preserve_first_run(sub_shape,
            "Copy this format for a complete, actionable Hands-Off (Concierge) request to COBE Marketing")

    # ─── Slide 23 (index 22) — Lead Times ────────────────────────────────────
    print("  Slide 23: Lead Times - updating for appendix move")
    s23 = slides[22]

    # Update title to indicate appendix
    lt_title = find_shape(s23, 'Google Shape;491;p22')
    if lt_title:
        set_text_preserve_first_run(lt_title, "Appendix: Lead Times Reference")

    # Update subtitle to add Hands-Off vs Hands-On note
    lt_sub = find_shape(s23, 'Google Shape;493;p22')
    if lt_sub:
        set_text_preserve_first_run(lt_sub,
            "Hands-Off: add time for COBE Marketing to create copy/design. Hands-On: faster if assets are already on-spec and final. Rush requests require director approval.")

    # ─── Slide 29 (index 28) — QR Code Tracking ──────────────────────────────
    print("  Slide 29: QR Code Tracking")
    s29 = slides[28]

    # Update the "Why Track" section - remove "central landing page" assumption
    why_body = find_shape(s29, 'Google Shape;569;p28')
    if why_body:
        set_text_preserve_first_run(why_body,
            "Print materials, posters, and classroom slides drive traffic\u2014but only if you measure it. Trackable QR codes let you see which placements actually generate scans and conversions.")

    # Update Implementation Steps
    impl_body = find_shape(s29, 'Google Shape;574;p28')
    if impl_body and impl_body.has_text_frame:
        tf = impl_body.text_frame
        steps = [
            "1. Choose your destination link (registration page, RSVP form, Boise State web page). If you don\u2019t have one, request a landing page via COBE Marketing.",
            "2. Add UTM parameters for digital tracking (if needed)",
            "3. Generate a unique QR code per placement (poster vs. slide vs. flyer)",
            "4. Name consistently: RHM-ShuttleStop-Feb25",
            "5. Review scan/traffic weekly in Flowcode dashboard",
        ]
        # Clear existing and re-add
        for i, p in enumerate(tf.paragraphs):
            p.clear()

        for i, step in enumerate(steps):
            if i == 0:
                p = tf.paragraphs[0]
            elif i < len(tf.paragraphs):
                p = tf.paragraphs[i]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = step
            run.font.size = Pt(8)
            run.font.name = 'Calibri'
            run.font.color.rgb = DARK_GRAY
            p.space_after = Pt(3)

    # Update the "Link all codes to your central landing page" reference
    # in the Flowcode recommendation section
    fc_body = find_shape(s29, 'Google Shape;571;p28')
    if fc_body:
        set_text_preserve_first_run(fc_body,
            "Free tier available. Create unique codes for each placement. Link each to your destination URL.")

    # ─── Slide 33 (index 32) — Key Insights ──────────────────────────────────
    print("  Slide 33: Key Insights - fixing 'we' language")
    s33 = slides[32]
    s33_sub = find_shape(s33, 'Google Shape;646;p32')
    if s33_sub:
        set_text_preserve_first_run(s33_sub, "What worked and what COBE Marketing learned")

    # ─── Slide 34 (index 33) — Section 08 divider ────────────────────────────
    print("  Slide 34: Section 08 divider - updating subtitle")
    s34 = slides[33]
    s34_sub = find_shape(s34, 'Google Shape;668;p33')
    if s34_sub:
        set_text_preserve_first_run(s34_sub, "Templates & Links Hub \u2022 FAQ \u2022 Contact Directory")

    # ─── Slide 35 (index 34) — FAQ ───────────────────────────────────────────
    print("  Slide 35: FAQ - adding template question, fixing language")
    s35 = slides[34]

    # Fix "we create them" → "COBE Marketing creates them" in FAQ
    faq_graphics_a = find_shape(s35, 'Google Shape;683;p34')
    if faq_graphics_a:
        set_text_preserve_first_run(faq_graphics_a,
            "A: No. Hands-Off: COBE Marketing creates them. Hands-On: use COBE templates from the Templates & Links Hub.")

    # Fix "how fast can you turn..." answer
    faq_speed_a = find_shape(s35, 'Google Shape;679;p34')
    if faq_speed_a:
        set_text_preserve_first_run(faq_speed_a,
            "A: Standard is 10 business days. Rush requests need director approval. See each channel slide for specific lead times.")

    # Fix Kate Watson reference
    faq_lastmin_a = find_shape(s35, 'Google Shape;687;p34')
    if faq_lastmin_a:
        set_text_preserve_first_run(faq_lastmin_a,
            "A: Contact Kate Watson directly. COBE Marketing will do its best to accommodate.")

    # Add new FAQ entries by modifying the bottom of the slide
    # Current last Q is at y=4206240, A at y=4480560.
    # The slide has room below or I can resize existing items.
    # Let me add template FAQ by adding textboxes
    add_multi_line_textbox(s35, Emu(457200), Emu(4750000), Emu(8229600), Emu(400000), [
        {"text": "Q: Where do I find the templates?", "size": 10, "bold": False, "color": DARK_GRAY},
        {"text": "A: See the Templates & Links Hub slide or the /templates folder in the playbook repository.",
         "size": 10, "bold": False, "color": DARK_GRAY},
    ])

    # ─── Slide 37 (index 36) — Let's Go Broncos ──────────────────────────────
    print("  Slide 37: Closing slide - fixing 'let's build' language")
    s37 = slides[36]
    closing = find_shape(s37, 'Google Shape;727;p36')
    if closing:
        set_text_preserve_first_run(closing,
            "Pick your path, send your request, and COBE Marketing is here to help you succeed.")

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 3: Structural changes
    # ═══════════════════════════════════════════════════════════════════════════
    print("\n--- Phase 3: Structural changes ---")

    # 3a. Create Templates & Links Hub slide (insert after slide 34 / index 33)
    print("  Creating Templates & Links Hub slide")

    # We need to add a new slide. Use a blank layout.
    # In python-pptx, to add a slide we need a layout.
    # Let's use the first layout (blank/default)
    slide_layout = prs.slide_layouts[0]  # Usually blank

    # Insert new slide - python-pptx appends by default
    # We'll add it at the end and then move it
    new_slide = prs.slides.add_slide(slide_layout)

    # Build Templates & Links Hub content
    # Title
    add_textbox(new_slide, Emu(457200), Emu(274320), Emu(8229600), Emu(548640),
                "Templates & Links Hub", font_size=28, bold=True, color=DARK_GRAY,
                font_name='Calibri')

    # Accent bar
    bar = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(457200), Emu(777240), Emu(1371600), Emu(54864)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BOISE_BLUE
    bar.line.fill.background()

    # Subtitle
    add_textbox(new_slide, Emu(457200), Emu(914400), Emu(8229600), Emu(274320),
                "Everything you need to get started - templates, tools, and reference links",
                font_size=12, color=MED_GRAY, font_name='Calibri')

    # Templates grid - 2 columns, 4 rows
    templates = [
        ("Creative Brief Template", "templates/creative-brief-template.md",
         "Submit campaign details to COBE Marketing"),
        ("Copy Prompts Template", "templates/copy-prompts-template.md",
         "Pre-written copy frameworks for each channel"),
        ("Campaign Calendar Template", "templates/campaign-calendar-template.md",
         "Plan multi-channel event timelines"),
        ("Budget Request Template", "templates/budget-request-template.md",
         "Estimate costs and request paid advertising"),
        ("Social Media Calendar Template", "Social Media Calendar Template.xlsx",
         "Weekly/monthly posting schedule"),
        ("Social Media Calendar Example", "Social Media Calendar Example.xlsx",
         "Sample calendar showing structure and pacing"),
    ]

    col_w = Emu(3931920)
    row_h = Emu(500000)
    start_y = Emu(1300000)
    start_x = Emu(457200)
    gap_x = Emu(365760)
    gap_y = Emu(80000)

    for idx, (name, path, desc) in enumerate(templates):
        col = idx % 2
        row = idx // 2
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)

        box = new_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, row_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        box.line.width = Pt(0.75)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)

        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = name
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = BOISE_BLUE
        r1.font.name = 'Calibri'

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(8)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = 'Calibri'

        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = f"\u2192 {path}"
        r3.font.size = Pt(7)
        r3.font.color.rgb = MED_GRAY
        r3.font.name = 'Calibri'
        r3.font.italic = True

    # Add key links section at bottom
    add_multi_line_textbox(new_slide, Emu(457200), Emu(4100000), Emu(8229600), Emu(900000), [
        {"text": "KEY LINKS", "size": 10, "bold": True, "color": DARK_GRAY, "space_after": 3},
        {"text": "\u2022 All requests: cobe-info@boisestate.edu  |  Subject: PROGRAM REQUEST - Channel - Date Range",
         "size": 8, "color": DARK_GRAY, "space_after": 1},
        {"text": "\u2022 University calendar: boisestate.enterprise.localist.com  |  Brand guide: brandguide.boisestate.edu",
         "size": 8, "color": DARK_GRAY, "space_after": 1},
        {"text": "\u2022 Career Services: cobecareers@boisestate.edu  |  QR codes: flowcode.com",
         "size": 8, "color": DARK_GRAY, "space_after": 1},
    ])

    # Now move this slide from the end to after slide 34 (index 33)
    # The new slide was appended at the end (index = total-1)
    # We need to move it to index 34 (after index 33)
    from lxml import etree
    sldIdLst = prs.slides._sldIdLst
    # The new slide's sldId is the last one
    new_sldId = sldIdLst[-1]
    # Remove from current position
    sldIdLst.remove(new_sldId)
    # Insert at position 34 (after index 33 = Section 08 divider)
    sldIdLst.insert(34, new_sldId)
    print("  Templates & Links Hub slide created and positioned")

    # 3b. Move lead times slide (originally index 22) to appendix
    # After our insert above, slide indices have shifted.
    # Original slide 23 (index 22) is still at index 22.
    # We want to move it near the end, right before the FAQ slide.
    # Current order after Templates insert:
    #   ... index 33 = Section 08, index 34 = Templates Hub, index 35 = FAQ,
    #   index 36 = Team, index 37 = Closing
    # We want Lead Times after Templates Hub and before FAQ.
    # So move from index 22 to index 35.

    print("  Moving Lead Times to appendix position")
    lead_times_sldId = sldIdLst[22]
    sldIdLst.remove(lead_times_sldId)
    # After removal, everything shifts. Templates Hub is now at index 33.
    # FAQ is at index 34. We want to insert before FAQ, so at index 34.
    # Actually let's insert after Templates Hub (index 33), making it index 34.
    sldIdLst.insert(34, lead_times_sldId)

    # 3c. Delete Instagram duplicate (slide 12, originally index 11)
    # After the move above, slide 12 is still at index 11
    # (the moves above were of indices > 11)
    print("  Deleting Instagram duplicate slide (original slide 12)")
    delete_slide(prs, 11)

    # Also update Section 04 divider subtitle to remove "Lead Times" since it moved
    # Section 04 was originally slide 20 (index 19). After deleting slide 12 (index 11),
    # it's now at index 18.
    slides_after = list(prs.slides)
    for s in slides_after:
        for shape in s.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "Colors \u2022 Sizes \u2022 Lead Times \u2022 Non-Negotiables":
                set_text_preserve_first_run(shape, "Colors \u2022 Sizes \u2022 Non-Negotiables")
                print("  Updated Section 04 subtitle (removed Lead Times reference)")
                break

    # ═══════════════════════════════════════════════════════════════════════════
    # PHASE 4: Final cleanup
    # ═══════════════════════════════════════════════════════════════════════════
    print("\n--- Phase 4: Final cleanup ---")

    # Do one more pass for any remaining "us/we" patterns
    final_replacements = [
        ("Send us ", "Email COBE Marketing "),
        ("send us ", "email COBE Marketing "),
        ("We'll ", "COBE Marketing will "),
        ("We will ", "COBE Marketing will "),
        ("we'll ", "COBE Marketing will "),
        ("We handle", "COBE Marketing handles"),
        ("we handle", "COBE Marketing handles"),
        ("We review", "COBE Marketing reviews"),
        ("we review", "COBE Marketing reviews"),
        ("We create", "COBE Marketing creates"),
        ("we create", "COBE Marketing creates"),
        ("our templates", "COBE templates"),
    ]

    for old, new in final_replacements:
        c = deck_wide_replace(prs, old, new)
        if c > 0:
            print(f"  Final pass: replaced '{old}' → {c} occurrences")

    # Verify no xxx@email remains
    found_xxx = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "xxx@email" in shape.text_frame.text:
                    found_xxx = True
                    print(f"  WARNING: Found 'xxx@email' in shape '{shape.name}'")
                    replace_in_all_runs(shape, "xxx@email", "cobe-info@boisestate.edu")
    if not found_xxx:
        print("  ✓ No placeholder emails (xxx@email) found")

    # ═══════════════════════════════════════════════════════════════════════════
    # Save
    # ═══════════════════════════════════════════════════════════════════════════
    print(f"\nSaving to {DST}...")
    prs.save(DST)
    print(f"Done! Saved {DST}")
    print(f"Final slide count: {len(prs.slides)}")

if __name__ == "__main__":
    main()
